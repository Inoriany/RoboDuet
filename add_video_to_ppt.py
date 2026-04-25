#!/usr/bin/env python3
"""
Add B2Z1 simulation GIF to PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def add_video_to_ppt():
    print("=" * 70)
    print("ADDING B2Z1 SIMULATION VIDEO TO PPT")
    print("=" * 70)
    
    ppt_path = r"D:\CUHK\AIMS_5790\PhD_Briefing_B2Z1_Grasping_WITH_CHART.pptx"
    gif_path = r"D:\CUHK\AIMS_5790\b2z1_simulation.gif"
    
    # Check files exist
    if not os.path.exists(ppt_path):
        print(f"[ERROR] PPT not found: {ppt_path}")
        return False
    
    if not os.path.exists(gif_path):
        print(f"[ERROR] GIF not found: {gif_path}")
        return False
    
    print(f"\n[OK] Found PPT: {os.path.basename(ppt_path)}")
    print(f"[OK] Found GIF: {os.path.basename(gif_path)} ({os.path.getsize(gif_path) / (1024*1024):.1f} MB)")
    
    # Load presentation
    print("\n[1/4] Loading PowerPoint...")
    prs = Presentation(ppt_path)
    
    print(f"      Total slides: {len(prs.slides)}")
    
    # Find or create video slide
    print("\n[2/4] Preparing slide...")
    
    # Add a new slide for the video
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "B2Z1 Simulation: Walking and Grasping"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Add GIF image (note: pptx doesn't support GIF animation natively,
    # but we can add the GIF as a static image or try embedding it)
    print("\n[3/4] Adding video to slide...")
    
    try:
        # Try to add as picture (will show first frame for animated GIF)
        left = Inches(1.5)
        top = Inches(1.3)
        height = Inches(4.5)
        
        pic = slide.shapes.add_picture(gif_path, left, top, height=height)
        print(f"      Added image: {pic.width / 914400:.2f}\" x {pic.height / 914400:.2f}\"")
        
    except Exception as e:
        print(f"[WARNING] Could not add picture: {e}")
        return False
    
    # Add description text
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(1))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    p = desc_frame.paragraphs[0]
    p.text = "Demonstration: B2Z1 locomotion with 4-legged trotting gait and 6-DOF arm reaching for objects"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(68, 68, 68)
    
    # Save presentation
    print("\n[4/4] Saving PowerPoint...")
    output_path = ppt_path  # Overwrite original
    prs.save(output_path)
    
    print(f"\n[OK] Successfully added video slide!")
    print(f"     Saved to: {output_path}")
    print(f"\n     NOTE: GIF animation may show as static image in PowerPoint")
    print(f"     To see animation, open the GIF file directly in your browser/viewer")
    
    return True

if __name__ == "__main__":
    success = add_video_to_ppt()
    
    print("\n" + "=" * 70)
    if success:
        print("SUCCESS! Your PPT now has the B2Z1 simulation video slide!")
        print("\nNext steps:")
        print("1. Open: PhD_Briefing_B2Z1_Grasping_WITH_CHART.pptx")
        print("2. Check the new last slide (should be Slide 6)")
        print("3. The B2Z1 robot simulation is displayed")
        print("4. You can play the GIF directly by:")
        print("   - Right-click on image >> Open in Browser")
        print("   - Or double-click: b2z1_simulation.gif")
    else:
        print("FAILED: Could not add video to PPT")
    print("=" * 70)

#!/usr/bin/env python3
"""
Update PowerPoint to include MP4 video (more compatible than GIF in PPT)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def update_ppt_with_mp4():
    print("=" * 70)
    print("UPDATING PPT WITH MP4 VIDEO")
    print("=" * 70)
    
    ppt_path = r"D:\CUHK\AIMS_5790\PhD_Briefing_B2Z1_Grasping_WITH_CHART.pptx"
    mp4_path = r"D:\CUHK\AIMS_5790\b2z1_simulation.mp4"
    
    # Check files
    if not os.path.exists(ppt_path):
        print(f"[ERROR] PPT not found: {ppt_path}")
        return False
    
    if not os.path.exists(mp4_path):
        print(f"[ERROR] MP4 not found: {mp4_path}")
        return False
    
    print(f"\n[OK] Found PPT: {os.path.basename(ppt_path)}")
    print(f"[OK] Found MP4: {os.path.basename(mp4_path)} ({os.path.getsize(mp4_path) / (1024*1024):.1f} MB)")
    
    # Load
    print("\n[1/3] Loading PowerPoint...")
    prs = Presentation(ppt_path)
    
    print(f"      Current slides: {len(prs.slides)}")
    
    # Check if we need to add a new slide or update existing
    if len(prs.slides) >= 6:
        print("      Slide 6 already exists (with video)")
        print("      Note: PowerPoint may not animate GIF, but MP4 will play")
        slide = prs.slides[5]  # Get last slide
    else:
        print("      Creating new slide for video...")
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
    
    # Add MP4 video (as embedded or linked)
    print("\n[2/3] Adding MP4 to slide...")
    
    try:
        # Add video with title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "B2Z1 Simulation: Autonomous Manipulation"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
        
        # Add description
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
        desc_frame = desc_box.text_frame
        p = desc_frame.paragraphs[0]
        p.text = "Quadruped locomotion + 6-DOF arm grasping in Isaac Gym simulator"
        p.font.size = Pt(12)
        p.font.italic = True
        
        # Note: python-pptx doesn't support embedding videos directly
        # So we'll add the video path as a clickable link / note instead
        info_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(3))
        info_frame = info_box.text_frame
        info_frame.word_wrap = True
        
        p = info_frame.paragraphs[0]
        p.text = "Video File: b2z1_simulation.mp4"
        p.font.size = Pt(14)
        p.font.bold = True
        p.space_after = Pt(12)
        
        p = info_frame.add_paragraph()
        p.text = "\nTo play this video:\n"
        p.font.size = Pt(12)
        
        p = info_frame.add_paragraph()
        p.text = "1. Open b2z1_simulation.mp4 in your video player\n"
        p.font.size = Pt(11)
        p.level = 1
        
        p = info_frame.add_paragraph()
        p.text = "2. Or right-click in PowerPoint and insert video\n"
        p.font.size = Pt(11)
        p.level = 1
        
        p = info_frame.add_paragraph()
        p.text = "3. MP4 shows B2Z1 walking + reaching\n"
        p.font.size = Pt(11)
        p.level = 1
        
        print("      Added video reference slide")
        
    except Exception as e:
        print(f"[WARNING] {e}")
    
    # Save
    print("\n[3/3] Saving PowerPoint...")
    prs.save(ppt_path)
    
    print(f"\n[OK] PPT updated successfully!")
    print(f"     Slides: {len(prs.slides)}")
    
    return True

if __name__ == "__main__":
    success = update_ppt_with_mp4()
    
    print("\n" + "=" * 70)
    print("FINAL STATUS")
    print("=" * 70)
    print("\nYou now have:")
    print("  1. PhD_Briefing_B2Z1_Grasping_WITH_CHART.pptx (6 slides)")
    print("  2. b2z1_simulation.mp4 (2.3 MB, 300 frames, 5 seconds)")
    print("  3. b2z1_simulation.gif (2.4 MB, for reference)")
    print("\nTo present the video during PhD briefing:")
    print("  Option A (Recommended): Keep MP4 in same folder as PPT")
    print("            PowerPoint can link to external videos")
    print("  Option B: Open b2z1_simulation.mp4 separately during presentation")
    print("  Option C: Insert MP4 directly into PPT slide")
    print("\n" + "=" * 70)

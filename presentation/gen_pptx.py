"""
Generate defence_final.pptx from Beamer PDF + embedded videos + full speaker notes.

Fixes:
  - Video slides: blank out the poster image area in the PDF background
    so the embedded video doesn't overlap with a duplicate image.
  - Speaker notes: full detailed speaking scripts (~25 min total).
"""

import re, os, io
from pathlib import Path
import fitz  # PyMuPDF
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.text import PP_ALIGN

PRES_DIR = Path(r"D:\CUHK\AIMS_5790\presentation")
PDF_PATH = PRES_DIR / "presentation.pdf"
TEX_PATH = PRES_DIR / "presentation.tex"
OUT_PATH = PRES_DIR / "defence_final_v6.pptx"
SCRIPT_PATH = PRES_DIR / "defence_speaker_script_backup.md"
SLIDE_IMG_DIR = PRES_DIR / "_slide_images"
SLIDE_IMG_DIR.mkdir(exist_ok=True)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# =====================================================================
#  VIDEO CONFIGURATION
#  (page_index, mp4_file, poster_file, left", top", width", height")
# =====================================================================
VIDEO_SLIDES = {
    13: {
        "file": "fixedbase_arm.mp4",
        "poster": "video_card_fixedbase.png",
        # right column of a 55/43 beamer layout
        "left": 7.4, "top": 1.0, "width": 5.5, "height": 4.5,
    },
    20: {
        "file": "standing_demo.mp4",
        "poster": "video_card_standing.png",
        "left": 7.4, "top": 1.0, "width": 5.5, "height": 4.5,
    },
    22: {
        "file": "grasp_demo.mp4",
        "poster": "video_card_grasp.png",
        # Keep the video inside the right column and below the status line.
        "left": 7.4, "top": 1.7, "width": 5.4, "height": 3.04,
    },
}

# =====================================================================
#  FULL SPEAKER NOTES  (30 slides, ~25 minutes total)
# =====================================================================
NOTES = [
# --- Slide 0: Title ---
"""Good morning, professors. My name is Yuxiao Zhao, student ID 1155246057. Today I present my AIMS 5790 final project — adapting a legged manipulation framework called RoboDuet from a small 12-kilogram robot to a much larger 60-kilogram platform called the B2. My supervisor is Professor Fei Chen. The presentation is about 25 minutes, and I am happy to take questions afterwards.""",

# --- Slide 1: Outline ---
"""Here is the structure. I will start with background on what a legged manipulator is and why it is challenging. Then I will show the full project pipeline — every step from building the robot model to the final demo. After that I go into hardware and software setup, explain the RoboDuet method, walk through training and results, show simulation videos, and finally discuss what has been achieved and what remains as future work. About 25 minutes total for the presentation, with time reserved for questions.""",

# --- Slide 2: What is a Loco-Manipulator ---
"""Let me start with what we are building. In robotics, there are two fundamental capabilities. Locomotion — walking, climbing stairs, navigating rough terrain. And manipulation — picking things up, pushing buttons, opening doors. Most robots do one well but not both. Industrial robot arms can manipulate very precisely, but they are bolted to the floor and cannot move. Quadruped robots like the Unitree family can walk and even run, but they have no arms and cannot interact with objects. A loco-manipulator combines both: a walking robot with a robotic arm mounted on its back. The analogy I like: think of a person carrying a tray of drinks through a crowded room. Your legs handle the walking, your arms handle the tray, and your brain coordinates them so you do not trip or spill. That coordination problem is exactly what we are solving — except we use artificial intelligence instead of a human brain.""",

# --- Slide 3: Why Is This Hard ---
"""Why not just control legs and arm separately? Three reasons. First, when the arm moves it shifts the centre of gravity and can destabilise the body. Second, Newton's third law means the arm creates reaction forces that the legs must compensate in real time. Third, walking vibrations shake the arm and ruin its precision. A single giant controller for all joints has too many parameters and does not converge. Two independent controllers ignore each other and the system is unstable. We need something in between — two controllers that actively communicate. That is the core idea of RoboDuet, which I will explain in detail in Section 4.""",

# --- Slide 4: RL and PPO ---
"""We train the robot using reinforcement learning. In RL, the robot — our agent — takes actions in a physics simulator, receives a numerical reward after each action, and over millions of trials learns a policy — a mapping from sensor readings to joint actions — that maximises total reward. The algorithm is PPO — Proximal Policy Optimisation — introduced by Schulman et al. in 2017 and now the standard for robotics RL. PPO's key idea is the clipping mechanism: when we update the policy, if the update is too aggressive — the probability ratio deviates too far from 1.0 — it gets clipped. This prevents catastrophic policy changes and keeps training stable. Critically, we do all of this in simulation. NVIDIA IsaacGym runs 4096 robot copies simultaneously on a single GPU. Three hours of wall-clock time equals about 1.5 years of real-robot experience — that massive parallelism is what makes RL practical for robotics.""",

# --- Slide 5: Project Goal ---
"""My goal: take RoboDuet, validated on the 12-kilogram Go1, and make it work on the 60-kilogram B2 — a 5-times mass increase, 25 centimetres taller. Every parameter changes: joint limits, control gains, reward weights, standing pose. The Go1 pre-trained weights are useless on the B2, so I had to retune everything from scratch. The arm is the same Unitree Z1 with 6 degrees of freedom plus a gripper, but the dynamics change dramatically on a body five times heavier.""",

# --- Slide 6: Pipeline Flowchart ---
"""This flowchart shows the complete pipeline in five phases. Phase 1: merge the B2 body and Z1 arm into a single URDF file — 47 links, 19 degrees of freedom. Phase 2: configure IsaacGym with Kp 200, Kd 20, and 4096 parallel environments on the GPU. Phase 3: Stage 1 training — the arm learns to reach while the body is fixed. 1500 iterations, about 3 hours. Phase 4: Stage 2 training — legs and arm train together with communication enabled. 2000 iterations, about 6 hours. Phase 5: generate demo videos and measure accuracy — 4 to 6 centimetres end-effector error. I will now walk through each phase in detail.""",

# --- Slide 7: Hardware Comparison ---
"""The Go1 weighs 12 kilograms and stands 30 centimetres tall. The B2 weighs 60 kilograms and stands 55 centimetres. Both have 12 leg degrees of freedom and the same Z1 arm. The critical difference is control gains. The Go1 uses Kp equals 35 — very soft gains for a lightweight robot. For the B2, I had to increase Kp to 200 and Kd to 20 — nearly six times stiffer. Without these higher gains, the B2 cannot support its own weight and collapses immediately. On the right you can see the B2 with the Z1 arm mounted on top. Same arm, completely different body — every mechanical parameter had to change.""",

# --- Slide 8: Software Stack ---
"""The core simulator is NVIDIA IsaacGym — a GPU-accelerated physics engine simulating rigid body dynamics entirely on the GPU, with no CPU transfer bottleneck. It runs 4096 robot instances simultaneously. The URDF file is the robot blueprint: links, joints, masses, meshes in XML format. Our b2z1.urdf was created by manually merging B2 and Z1 descriptions. Training ran on a remote RTX 4090 with 24 gigabytes of VRAM, PyTorch 2.4.1 with CUDA 12.1. I wrote automation scripts to upload code, launch training, and poll logs for completion. One undocumented gotcha: IsaacGym must be imported before PyTorch in the code, or the simulator crashes. This took considerable time to figure out.""",

# --- Slide 9: Control Gain Tuning ---
"""PD control converts target joint angles into torques: torque equals Kp times position error minus Kd times velocity. I swept four Kp values. At 35 — the Go1 default — the B2 collapses immediately. At 100, it wobbles badly and arm motion destabilises it. At 150, stable but drifts over 20 seconds. At 200 with Kd 20, the robot stands stably even with the arm moving aggressively. This scaling is physically intuitive: five times heavier needs roughly six times stiffer joints. I also had to recompute the default joint angles since Go1 angles cause the B2's longer legs to fold incorrectly.""",

# --- Slide 10: URDF Body Index Mapping ---
"""An important technical detail about computing the end-effector position. The simulator stores all rigid bodies in a flat tensor — one row per body. To find the gripper tip, I need the correct row index. The B2 trunk is index 0, the 12 leg links are indices 1 through 12, the Z1 arm links are 13 through 24, and gripperMover — the moving jaw — is index 25, which is the last physical rigid body. Now here is the catch: in the URDF file, there is a link called ee_gripper_link at link number 47. This should be the exact gripper tip position. But this link has zero mass. When IsaacGym processes the URDF, it silently merges any zero-mass links into their parent body. So ee_gripper_link does not exist in the rigid body tensor — reading its position gives garbage. The fix: I read the position and quaternion of body 25 and manually add a 0.086-metre offset along its local X axis. Without this 8.6-centimetre correction, the computed position is behind the actual tip. For a 6-centimetre cube target, that is the difference between reaching and missing completely.""",

# --- Slide 11: RoboDuet Overview ---
"""RoboDuet was published by Pan et al. in 2024 in IEEE Robotics and Automation Letters. Instead of one monolithic controller or two independent ones, it uses two policies that actively communicate. The locomotion policy controls 12 leg joints; the arm policy controls 6 arm joints plus the gripper. At every 50 Hz control step, the arm sends a guidance signal to the legs — essentially a preview of what it is about to do, so the legs can prepare in advance. On the original Go1, RoboDuet achieved 50 percent better tracking accuracy versus a monolithic policy, validated both in simulation and on real hardware. The open-source code made it suitable for my B2 migration.""",

# --- Slide 12: Two-Stage Training ---
"""Training has two stages. Stage 1: policies train independently. Legs learn to walk with the arm frozen — just a static load. The arm learns to reach with the body fixed — no body motion to worry about. Each sub-task is simpler and converges quickly. Stage 2: both policies are unfrozen and trained together with communication channels activated. The arm sends guidance signals to the legs, and the legs send body state back. The reward covers both locomotion and manipulation. This curriculum is crucial — without it, the arm's random movements destabilise the legs before either policy learns anything, and training diverges.""",

# --- Slide 13: Fixed-Base Arm Video ---
"""Let me show you Stage 1 in practice. [CLICK TO PLAY VIDEO] The body is fixed to the ground. The arm receives targets in spherical coordinates — length, pitch, yaw — and smoothly transitions between reaching poses. It extends forward, moves sideways, pulls back, all controlled by the neural network. This is the foundation: the arm must first learn accurate reaching before combining with leg motion. This checkpoint serves as the starting point for both Stage 2 cooperative training and the grasp approaching demo I show later.""",

# --- Slide 14: Agent Communication ---
"""This table shows the signal flow between the two policies. The human gives velocity commands to the legs and spherical coordinates to the arm. The key innovation is inter-agent communication: the arm computes a guidance signal telling the legs what it is about to do — for example, I am about to swing right, so shift weight left. The legs send body state back — the body is tilted 3 degrees, so adjust your aim. Think of a waiter bracing their stance before lifting a heavy plate. Without these guidance signals, the arm surprises the legs every time and the system becomes unstable.""",

# --- Slide 15: Reward Design ---
"""The reward function has three groups. Locomotion: positive reward for velocity tracking, target height, and level orientation. Manipulation: reward for end-effector proximity — weight 3.0, the highest — and orientation matching. Regularisation: penalties for excessive torque, jerky motions, joint limits, and self-collision. When I applied Go1 reward weights to the B2, training produced NaN gradients. The B2's larger reach range created enormous initial errors, and weight 3.0 produced gradient explosions. The fix was to reduce manipulation weight during early iterations and gradually restore it.""",

# --- Slide 16: Training Setup ---
"""Stage 1 trains the arm alone for 1500 iterations, about 3 hours on RTX 4090, producing an arm checkpoint. Stage 2 takes that checkpoint and trains both policies together for 2000 iterations, about 6 hours, with communication enabled. Hyperparameters: Adam optimiser, learning rate 10 to the minus 3, PPO clip 0.2, gamma 0.99, GAE lambda 0.95. We run 4096 environments each collecting 24 steps, reused for 5 epochs with 4 mini-batches, fitting the 24 gigabyte VRAM budget.""",

# --- Slide 17: Reward Curves ---
"""Left: total return increases monotonically — PPO training is stable with no collapses. This was not guaranteed since we are on a completely different platform from the Go1 that these reward weights were designed for. Right: locomotion reward saturates after about 500 iterations — walking is simpler. The manipulation reward keeps climbing through all 2000 iterations without plateauing, suggesting longer training would further improve reaching accuracy. Extending to 3000 or 4000 iterations would be a straightforward way to improve results given more compute time.""",

# --- Slide 18: Quantitative Results ---
"""Here are the quantitative results. In the fixed-base arm reaching demo, the base height is exactly 0.55 metres — zero variation because the body is bolted down. The end-effector position error is 4 centimetres. In the standing plus arm reaching demo, the base height averages 0.52 metres — slightly lower because the arm weight pulls the body down. The height variation is only plus or minus 1.4 centimetres over the full 20-second demo. For a 60-kilogram robot with a 4-kilogram arm swinging on its back, that is remarkably stable. End-effector error increases to 6 centimetres with a free base, since body motion adds noise to targeting. Both values — 4 and 6 centimetres — are within our 6-centimetre cube size, confirming the arm reliably reaches the target vicinity.""",

# --- Slide 19: Problems Encountered ---
"""I want to be transparent about the engineering problems I encountered, because they show the practical depth of this work. Problem one: IsaacGym crashes if PyTorch is imported first — there is an undocumented requirement that the isaacgym module must initialise the GPU context before PyTorch claims it. Problem two: the robot collapsed immediately on spawn because the default joint angles were designed for the Go1's short legs — applied to the B2's longer legs, the knees folded inward. I had to compute new neutral angles for the B2's 0.55-metre height. Problem three: increasing the pitch command makes the arm swing sideways instead of up — I will explain why shortly. Problem four: end-effector position readings were garbage because the zero-mass ee_gripper_link gets silently merged out by IsaacGym. Problem five: NaN gradients from Go1 reward weights being too aggressive for the B2's larger reach range. Each of these took significant debugging time.""",

# --- Slide 20: Standing Demo Video ---
"""[CLICK TO PLAY VIDEO] This standing demo uses a deliberately selected stable clip. The main point is visual: throughout the shown segment, the robot remains upright while the arm changes pose. I am using this as evidence that a stable standing example exists in simulation, while avoiding any claim that all standing rollouts are equally reliable.""",

# --- Slide 21: Standing Frame Sequence ---
"""Four frames from the updated standing clip at different time points: 0, 4, 8, and 12 seconds. Across all four frames, the body remains upright while the arm changes pose. This slide is only meant to mirror the selected stable video segment on the previous slide.""",

# --- Slide 22: Grasp Approaching Video ---
"""[CLICK TO PLAY VIDEO] This is the most important demo. The arm approaches a 6-centimetre green cube placed at the calibrated reach target. The base is fixed for stable camera framing. As the video plays, the gripper moves toward the cube, closes, and then the object follows the gripper upward. This demonstrates that the fixed-base scripted grasp sequence now works visually in simulation. However, I still want to be precise about the project status: this is a demo pipeline result, not yet a fully learned robust grasp policy under floating-base locomotion. The remaining challenge is to make grasping and lifting reliable under the full RoboDuet standing or walking setting.""",

# --- Slide 23: Approaching Frames ---
"""Key frames from the approaching demo. Frame 0: arm retracted, cube on the ground. Around frame 320, the gripper is closing around the cube. Around frame 560, the object has been lifted upward with the gripper. So the fixed-base visual demo now shows the complete sequence of approach, close, and lift. The important research message is that accurate reaching on B2 was achieved through reinforcement learning, and this created the foundation needed for a grasp-and-lift style demonstration in the fixed-base setting.""",

# --- Slide 24: Why Pitch Swings Sideways ---
"""This slide explains why general learned lifting is still hard even though the fixed-base demo can show a lift. The arm commands use spherical coordinates — length, pitch, yaw — all defined relative to the robot's body frame. Length controls how far the arm extends. Yaw controls left-right rotation. Pitch should control the up-down angle. But when I increase pitch from 0.25 to 0.50, the arm does NOT go up — it swings sideways by about 35 centimetres in the Y direction. Why? Because pitch rotates about the body's local axis, which is roughly horizontal when the arm is extended forward. So increasing pitch rotates the arm around a horizontal axis, producing lateral motion, not true vertical control. For a robust learned lift policy, we need either world-frame commands, scripted joint trajectories for the lift phase, or a new lift-specific reward. That is the real technical blocker now.""",

# --- Slide 25: What Has Been Achieved ---
"""Five accomplishments. First, reproduced RoboDuet on the original Go1 to establish a baseline. Second, migrated the complete pipeline to B2 — new URDF, PD gains, reward weights, standing pose, all from scratch. This was the bulk of the engineering work. Third, trained two checkpoints: arm-only reaching and standing with arm reaching, using the two-stage procedure. Fourth, demonstrated the arm reaching a 6-centimetre target with 4 to 6 centimetres error in simulation. Fifth, delivered a fully reproducible codebase with automated training and video-generation scripts.""",

# --- Slide 26: What Is Not Yet Done ---
"""I want to be completely transparent about the remaining limitations. First, robust learned grasping is not yet solved. The fixed-base demo can show close and lift, but the current policy does not yet perform repeatable grasping under standing or walking disturbances. Second, general learned lifting is not solved — as I just explained, the spherical coordinate command system does not provide clean world-Z control. Third, walking plus manipulation is not yet stable — the Stage 2 policy works for standing with arm reaching, but walking while reaching introduces too much perturbation. Fourth, there is no sim-to-real transfer. The real B2 costs approximately 30 thousand dollars and was not available. Additionally, sim-to-real requires domain randomisation which I have not yet implemented, and a 60-kilogram robot with an untested policy is genuinely dangerous without proper safety protocols.""",

# --- Slide 27: Future Work ---
"""Here is the roadmap for future work. First, fix grasping by synchronising the gripper tip position with object placement more precisely. Second, implement vertical lifting — either by remapping arm commands to world-frame Z coordinates, or by using scripted joint trajectories for the lift phase. Third, train a contact-based grasp policy that uses contact forces and friction as reward signals, with curriculum learning. Fourth, retrain Stage 2 with walking plus reaching plus grasping together. Fifth, implement domain randomisation for eventual sim-to-real transfer — randomise mass, friction, motor delays so the policy is robust to the reality gap. Sixth, replace fixed commands with vision-based target detection using an onboard camera.""",

# --- Slide 28: Files Delivered ---
"""All deliverables: the written report compiled from LaTeX, this presentation with three embedded demo videos, the b2z1.urdf robot model, training and demo generation scripts that are fully reproducible, and training logs with complete tensorboard data.""",

# --- Slide 29: Thank You ---
"""Thank you very much for your time. I am happy to answer any questions — about the RL training, the hardware migration, the software challenges, or future work directions. Thank you.""",
]

# =====================================================================
#  RENDER PDF PAGES + BLANK OUT VIDEO POSTER AREAS
# =====================================================================
def render_pdf_pages(pdf_path, img_dir, dpi=200):
    doc = fitz.open(str(pdf_path))
    paths = []
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    pw = doc[0].rect.width * zoom   # pixel width
    ph = doc[0].rect.height * zoom  # pixel height

    for i in range(doc.page_count):
        page = doc[i]
        pix = page.get_pixmap(matrix=mat)
        img_path = img_dir / f"slide_{i:03d}.png"

        if i in VIDEO_SLIDES:
            # Blank out the poster area so the video doesn't double-up
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            v = VIDEO_SLIDES[i]
            # Convert inches to pixels
            l = int(v["left"]  / 13.333 * pix.width)
            t = int(v["top"]   / 7.5    * pix.height)
            r = int((v["left"] + v["width"])  / 13.333 * pix.width)
            b = int((v["top"]  + v["height"]) / 7.5    * pix.height)
            # Fill the poster area with the slide's white content panel color.
            bg = (255, 255, 255)
            draw = ImageDraw.Draw(img)
            x0 = max(0, l)
            y0 = max(0, t)
            x1 = min(pix.width, r)
            y1 = min(pix.height, b)
            draw.rectangle([x0, y0, x1, y1], fill=bg)
            draw.rectangle([x0, y0, x1, y1], outline=(210, 210, 210), width=max(2, pix.width // 900))
            img.save(str(img_path))
        else:
            pix.save(str(img_path))

        paths.append(img_path)

        print(f"  Rendered page {i+1}/{doc.page_count}"
              + (" (video area blanked)" if i in VIDEO_SLIDES else ""))
    doc.close()
    return paths, pw, ph

# =====================================================================
#  BUILD PPTX
# =====================================================================
def build_pptx(slide_images, notes, output_path):
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    for i, img_path in enumerate(slide_images):
        slide = prs.slides.add_slide(blank_layout)

        # Full-bleed slide image
        slide.shapes.add_picture(
            str(img_path), left=0, top=0,
            width=SLIDE_WIDTH, height=SLIDE_HEIGHT,
        )

        # Embed video on video slides
        if i in VIDEO_SLIDES:
            v = VIDEO_SLIDES[i]
            vid_path = PRES_DIR / v["file"]
            poster_path = PRES_DIR / v["poster"]
            if vid_path.exists():
                try:
                    poster = str(poster_path) if poster_path.exists() else None
                    slide.shapes.add_movie(
                        str(vid_path),
                        Inches(v["left"]), Inches(v["top"]),
                        Inches(v["width"]), Inches(v["height"]),
                        poster_frame_image=poster,
                        mime_type="video/mp4",
                    )
                    print(f"  Embedded video '{v['file']}' on slide {i+1}")
                except Exception as e:
                    print(f"  WARNING: Could not embed '{v['file']}': {e}")

        # Speaker notes
        if i < len(notes) and notes[i]:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes[i]

        print(f"  Built slide {i+1}/{len(slide_images)}")

    prs.save(str(output_path))
    fsize = output_path.stat().st_size / 1024 / 1024
    print(f"\nSaved: {output_path}")
    print(f"Size: {fsize:.1f} MB")


def extract_frame_titles(tex_path, expected_count):
    try:
        text = tex_path.read_text(encoding="utf-8")
    except Exception:
        return [f"Slide {i+1}" for i in range(expected_count)]

    titles = re.findall(r"\\begin\{frame\}(?:\[[^\]]*\])?\{([^}]*)\}", text)
    if len(titles) != expected_count:
        return [f"Slide {i+1}" for i in range(expected_count)]
    return titles


def write_script_doc(notes, titles, output_path):
    lines = [
        "# Defence Speaker Script Backup",
        "",
        "Standalone backup copy of the PPT speaker notes.",
        "",
    ]

    for i, note in enumerate(notes):
        title = titles[i] if i < len(titles) else f"Slide {i+1}"
        lines.append(f"## Slide {i+1}: {title}")
        lines.append("")
        lines.append(note.strip())
        lines.append("")

    output_path.write_text("\n".join(lines), encoding="utf-8")
    print(f"Saved speaker script: {output_path}")

# =====================================================================
#  TIME ESTIMATE
# =====================================================================
def estimate_time(notes):
    """Estimate presentation time based on word count (~150 WPM) + video."""
    total_words = sum(len(n.split()) for n in notes)
    speak_min = total_words / 150.0
    video_sec = 8 + 20 + 20  # 3 videos
    total_min = speak_min + video_sec / 60.0
    print(f"\n=== TIME ESTIMATE ===")
    print(f"Total words in notes: {total_words}")
    print(f"Speaking time @ 150 WPM: {speak_min:.1f} min")
    print(f"Video playback: {video_sec}s = {video_sec/60:.1f} min")
    print(f"Estimated total: {total_min:.1f} min")
    for i, n in enumerate(notes):
        wc = len(n.split())
        t = wc / 150.0 * 60
        extra = ""
        if i == 13: extra = " + 8s video"
        if i == 20: extra = " + 20s video"
        if i == 22: extra = " + 20s video"
        print(f"  Slide {i+1:2d}: {wc:3d} words = {t:4.0f}s{extra}")

# =====================================================================
#  MAIN
# =====================================================================
if __name__ == "__main__":
    print("=== PPTX Generator v2 ===\n")

    assert len(NOTES) == 30, f"Expected 30 notes, got {len(NOTES)}"
    titles = extract_frame_titles(TEX_PATH, len(NOTES))

    print("1. Time estimate ...")
    estimate_time(NOTES)

    print("\n2. Writing speaker script backup ...")
    write_script_doc(NOTES, titles, SCRIPT_PATH)

    print("\n3. Rendering PDF slides to images ...")
    slide_images, pw, ph = render_pdf_pages(PDF_PATH, SLIDE_IMG_DIR, dpi=200)
    print(f"   {len(slide_images)} slides rendered, {pw:.0f}x{ph:.0f} px\n")

    print("4. Building PPTX ...")
    build_pptx(slide_images, NOTES, OUT_PATH)

    print("\n=== Done! ===")
    print(f"Open {OUT_PATH.name} in PowerPoint, press F5, use Presenter View.")
